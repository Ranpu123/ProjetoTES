import os
from pptx import Presentation
from spire.presentation import Presentation as pre
from spire.presentation import FileFormat
from spire.presentation.common import *
import hashlib
from datetime import datetime

TEMPLATE_PATH = "projetoApp\certificate_template\modelo.pptx"

def generateCertificado(nomeCompleto, temaEvento, dataEvento, cargaHoraria):   
    certificado = Presentation(TEMPLATE_PATH)
    replace_text_in_presentation(certificado, "FULLNAME", nomeCompleto)
    replace_text_in_presentation(certificado, "EVENT", temaEvento)
    replace_text_in_presentation(certificado, "DATE", dataEvento)
    replace_text_in_presentation(certificado, "TIME", cargaHoraria)

    #THIS IS FINE
    current_timestamp = datetime.now().timestamp()
    fname = nomeCompleto+temaEvento+str(int(current_timestamp))
    fname = hashlib.md5(bytes(fname, "UTF-8"))
    hashString = fname.hexdigest()
    certificado.save("projetoApp/temp_certificate/"+hashString+".pptx")
    print(hashString)

    temp = pre()
    temp.LoadFromFile("projetoApp/temp_certificate/"+hashString+".pptx")
    temp.SaveToFile("projetoApp/pdf_out/"+hashString+".pdf", FileFormat.PDF)
    temp.Dispose()

    os.remove("projetoApp/temp_certificate/"+hashString+".pptx")

    return os.path.abspath("projetoApp/pdf_out/"+hashString+".pdf")

def replace_text_in_shapes(shapes, target_text, replacement_text):
    for shape in shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = run.text.replace(target_text, replacement_text)

def replace_text_in_presentation(presentation, target_text, replacement_text):
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                replace_text_in_shapes([shape], target_text, replacement_text)

